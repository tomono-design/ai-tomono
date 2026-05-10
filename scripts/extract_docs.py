"""
AI友野 - 資料テキスト抽出スクリプト
iCloud Drive の講演資料・出版物からテキストを抽出し、chunks.json を生成する
"""

import json
import os
import re
from pathlib import Path

# ---- 設定 ----
ICLOUD = Path(r"C:\Users\user\iCloudDrive")

SOURCE_DIRS = [
    ICLOUD / "研修講演資料",
    ICLOUD / "出版",
    ICLOUD / "あんど" / "講演" / "原稿",
    ICLOUD / "あんど" / "原稿",
    ICLOUD / "作品",
    ICLOUD / "NotebookLM" / "書籍",
]

OUTPUT_FILE = Path(__file__).parent.parent / "src" / "data" / "chunks.json"
CHUNK_SIZE = 600   # 文字数
OVERLAP = 100      # 重複文字数


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)
    except Exception as e:
        print(f"  PPTX エラー: {path.name} - {e}")
        return ""


def extract_pdf(path: Path) -> str:
    try:
        import fitz
        doc = fitz.open(str(path))
        texts = []
        for page in doc:
            texts.append(page.get_text())
        doc.close()
        return "\n".join(texts)
    except Exception as e:
        print(f"  PDF エラー: {path.name} - {e}")
        return ""


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        texts = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(texts)
    except Exception as e:
        print(f"  DOCX エラー: {path.name} - {e}")
        return ""


def split_chunks(text: str, source: str) -> list[dict]:
    text = re.sub(r'\s+', ' ', text).strip()
    if len(text) < 50:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + CHUNK_SIZE, len(text))
        chunk = text[start:end]
        if len(chunk) > 50:
            chunks.append({"text": chunk, "source": source})
        start += CHUNK_SIZE - OVERLAP
    return chunks


def main():
    all_chunks = []
    seen_files = set()

    OUTPUT_FILE.parent.mkdir(parents=True, exist_ok=True)

    for source_dir in SOURCE_DIRS:
        if not source_dir.exists():
            print(f"スキップ（存在しない）: {source_dir}")
            continue

        print(f"\n処理中: {source_dir.name}")
        files = list(source_dir.rglob("*.pptx")) + \
                list(source_dir.rglob("*.pdf")) + \
                list(source_dir.rglob("*.docx"))

        for f in files:
            if f.name in seen_files:
                continue
            seen_files.add(f.name)

            # 除外フォルダチェック
            skip_keywords = ["給与", "訴訟", "労務", "相談案件", "支援報告", "領収書", "請求書"]
            if any(kw in str(f) for kw in skip_keywords):
                continue

            print(f"  {f.name}")
            ext = f.suffix.lower()
            if ext == ".pptx":
                text = extract_pptx(f)
            elif ext == ".pdf":
                text = extract_pdf(f)
            elif ext == ".docx":
                text = extract_docx(f)
            else:
                continue

            source_label = f"{source_dir.name}/{f.name}"
            chunks = split_chunks(text, source_label)
            all_chunks.extend(chunks)

    print(f"\n合計チャンク数: {len(all_chunks)}")
    with open(OUTPUT_FILE, "w", encoding="utf-8") as fp:
        json.dump(all_chunks, fp, ensure_ascii=False, indent=2)
    print(f"保存完了: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
